from pathlib import Path

from pptx import Presentation


def extract_pptx(path: str | Path) -> dict:
    path = Path(path)

    presentation = Presentation(path)

    slides = []

    for slide_number, slide in enumerate(
        presentation.slides,
        start=1,
    ):
        parts = []

        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue

            value = shape.text.strip()

            if value:
                parts.append(value)

        slide_text = "\n".join(parts)

        if slide_text:
            slides.append(
                {
                    "slide_number": slide_number,
                    "text": slide_text,
                }
            )

    full_text = "\n\n".join(
        f"Slide {slide['slide_number']}\n{slide['text']}"
        for slide in slides
    )

    return {
        "text": full_text,
        "title": path.stem,
        "metadata": {
            "slide_count": len(presentation.slides),
            "slides": slides,
        },
    }
